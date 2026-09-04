"""L6′ 用の最小 Office フィクスチャを決定的に生成する。

`fixtures/eval/excel_ja`／`fixtures/eval/office_ja` の実データ棚卸し
（`benchmarks/deprecation_inventory.json`・生成は `generate_inventory.py`）で見つからなかった
「廃止」構造マーカーを、ここで最小サイズの Office ファイルとして補う。

再生成:
    .venv/bin/python fixtures/eval/deprecation_markers/generate_fixtures.py

非決定要素（生成時刻・乱数）を持ち込まない: 文書プロパティの created/modified は固定値、
シェイプの座標・テキストはすべてリテラル。既存の `fixtures/eval/excel_ja/assets/廃止スタンプ.png`
を読み込んで xlsx の画像覆いに使う（アセットの複製はしない）。
"""
from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path

from docx import Document
from docx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Font
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu

_FIXED_TIMESTAMP = datetime(2020, 1, 1, tzinfo=timezone.utc)
_AUTHOR = "sherpa-fixture-generator"

_ROOT = Path(__file__).resolve().parents[3]
_STAMP_PNG = _ROOT / "fixtures" / "eval" / "excel_ja" / "assets" / "廃止スタンプ.png"
_INPUTS_DIR = Path(__file__).resolve().parent / "inputs"


def _build_xlsx() -> None:
    """xlsx: 取り消し線・画像による覆い・非表示シート・非表示行/列を1ファイルに集約する。

    - B2（`廃止予定`）: フォントの取り消し線（`font.strike`）。
    - B3: `廃止スタンプ.png` をセル上に配置（画像による覆い）。台帳上の値は「使用中」のまま
      （断定はしない設計・`status_semantics.covered_sentence` と同じ思想）。
    - 4行目: 行非表示（`row_dimensions[4].hidden`）。
    - C列（`内部コード`）: 列非表示（`column_dimensions["C"].hidden`）。
    - シート「旧版」: シート非表示（`state="hidden"`）。
    - シート「内部退避」: 完全非表示（`state="veryHidden"`）。
    """
    wb = Workbook()
    wb.properties.creator = _AUTHOR
    wb.properties.created = _FIXED_TIMESTAMP
    wb.properties.modified = _FIXED_TIMESTAMP
    wb.properties.lastModifiedBy = _AUTHOR

    ws = wb.active
    ws.title = "対象"
    ws.append(["項目ID", "状態", "内部コード"])
    ws.append(["ITEM01", "廃止予定", "X01"])
    ws.append(["ITEM02", "使用中", "X02"])
    ws.append(["ITEM03", "非表示行", "X03"])
    ws["B2"].font = Font(strike=True)
    img = XLImage(str(_STAMP_PNG))
    ws.add_image(img, "B3")
    ws.row_dimensions[4].hidden = True
    ws.column_dimensions["C"].hidden = True

    hidden_sheet = wb.create_sheet("旧版")
    hidden_sheet.append(["旧税率", 0.03])
    hidden_sheet.sheet_state = "hidden"

    very_hidden_sheet = wb.create_sheet("内部退避")
    very_hidden_sheet.append(["削除予定コード X99"])
    very_hidden_sheet.sheet_state = "veryHidden"

    _INPUTS_DIR.mkdir(parents=True, exist_ok=True)
    wb.save(_INPUTS_DIR / "DEP-XLSX-MARKERS.xlsx")


def _add_deleted_run(paragraph, text: str) -> None:
    """python-docx に無い `w:del`（変更履歴削除）ブロックを生の OOXML で挿入する。"""
    p_el = paragraph._p
    del_el = p_el.makeelement(
        qn("w:del"), {qn("w:id"): "1", qn("w:author"): _AUTHOR, qn("w:date"): "2020-01-01T00:00:00Z"})
    r_el = del_el.makeelement(qn("w:r"), {})
    del_text_el = r_el.makeelement(qn("w:delText"), {})
    del_text_el.text = text
    r_el.append(del_text_el)
    del_el.append(r_el)
    p_el.append(del_el)


def _build_docx() -> None:
    """docx: 取り消し線（単/二重）・隠し文字・変更履歴削除を1ファイルに集約する。"""
    doc = Document()
    doc.core_properties.author = _AUTHOR
    doc.core_properties.created = _FIXED_TIMESTAMP
    doc.core_properties.modified = _FIXED_TIMESTAMP
    doc.core_properties.last_modified_by = _AUTHOR

    p1 = doc.add_paragraph()
    p1.add_run("旧料金プラン（取り消し線）").font.strike = True

    p2 = doc.add_paragraph()
    p2.add_run("旧仕様（二重取り消し線）").font.double_strike = True

    p3 = doc.add_paragraph()
    p3.add_run("内部メモ：非公開（隠し文字）").font.hidden = True

    p4 = doc.add_paragraph()
    _add_deleted_run(p4, "旧契約条項（変更履歴で削除）")

    _INPUTS_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(_INPUTS_DIR / "DEP-DOCX-MARKERS.docx")


def _build_pptx() -> None:
    """pptx: 前面図形による覆い（occluded／covered_by_text）・スライド外配置・非表示スライドを集約する。

    スライド1: 無地塗りの前面シェイプ（テキスト無し）が背面テキストへ完全重畳＝`occluded`。
    スライド2: 前面テキストシェイプが背面テキストへ完全重畳＝`covered_by_text`。
    スライド3: シェイプをスライド範囲外に配置＝`off_slide`。
    スライド4: スライド自体を非表示（`show="0"`）＝`hidden_slide`。
    """
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    blank_layout = prs.slide_layouts[6]

    bbox = (Emu(500000), Emu(500000), Emu(3000000), Emu(500000))

    slide1 = prs.slides.add_slide(blank_layout)
    back1 = slide1.shapes.add_textbox(*bbox)
    back1.text_frame.text = "旧機能C: 提供終了予定"
    from pptx.enum.shapes import MSO_SHAPE
    front1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, *bbox)
    front1.fill.solid()
    front1.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    front1.line.fill.background()
    front1.text_frame.text = ""

    slide2 = prs.slides.add_slide(blank_layout)
    back2 = slide2.shapes.add_textbox(*bbox)
    back2.text_frame.text = "旧料金体系: 月額1000円"
    front2 = slide2.shapes.add_textbox(*bbox)
    front2.text_frame.text = "廃止"

    slide3 = prs.slides.add_slide(blank_layout)
    off_left = Emu(prs.slide_width + Emu(1000000))
    off_shape = slide3.shapes.add_textbox(off_left, Emu(500000), Emu(2000000), Emu(500000))
    off_shape.text_frame.text = "旧仕様メモ（画面外配置）"

    slide4 = prs.slides.add_slide(blank_layout)
    slide4.shapes.add_textbox(Emu(500000), Emu(500000), Emu(3000000), Emu(500000)).text_frame.text = \
        "旧業務フロー（非表示スライド）"
    slide4.element.set("show", "0")

    prs.core_properties.author = _AUTHOR
    prs.core_properties.created = _FIXED_TIMESTAMP
    prs.core_properties.modified = _FIXED_TIMESTAMP
    prs.core_properties.last_modified_by = _AUTHOR

    _INPUTS_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(_INPUTS_DIR / "DEP-PPTX-MARKERS.pptx")


def main() -> None:
    _build_xlsx()
    _build_docx()
    _build_pptx()
    print(f"generated 3 fixtures under {_INPUTS_DIR}")


if __name__ == "__main__":
    main()
